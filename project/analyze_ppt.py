from pptx import Presentation
from pptx.util import Inches, Pt
import os

def analyze_ppt(filepath, name):
    print("=" * 60)
    print(f"分析PPT: {name}")
    print("=" * 60)
    
    prs = Presentation(filepath)
    print(f"幻灯片数量: {len(prs.slides)}")
    print(f"幻灯片宽度: {prs.slide_width.inches} inches")
    print(f"幻灯片高度: {prs.slide_height.inches} inches")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- 幻灯片 {i+1} ---")
        print(f"布局: {slide.slide_layout.name if slide.slide_layout else 'Unknown'}")
        
        for shape in slide.shapes:
            print(f"  Shape: {shape.shape_type}, Name: {shape.name}")
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        print(f"    Text: {text[:100]}{'...' if len(text) > 100 else ''}")
            if shape.has_table:
                table = shape.table
                print(f"    Table: {table.rows.__len__()} rows x {len(table.columns)} cols")
                for row_idx, row in enumerate(table.rows):
                    if row_idx < 5:  # 只显示前5行
                        row_data = [cell.text[:20] for cell in row.cells]
                        print(f"      Row {row_idx}: {row_data}")

# 分析两个PPT
analyze_ppt(r'c:\Users\宇宙无敌智慧大将军\Desktop\test\data_analysis_agent\project\真味道&好汤面.pptx', '真味道&好汤面.pptx')
print("\n" * 2)
analyze_ppt(r'c:\Users\宇宙无敌智慧大将军\Desktop\test\data_analysis_agent\project\营运月会专案-香爆脆.pptx', '营运月会专案-香爆脆.pptx')
