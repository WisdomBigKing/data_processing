"""
报告生成器
支持生成Word和PPT格式的报告
"""

import os
from typing import Dict, List, Any, Optional
from dataclasses import dataclass
from abc import ABC, abstractmethod
from datetime import datetime

# Word生成
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE

# PPT生成
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

from .data_analyzer import AnalysisResult


@dataclass
class ReportConfig:
    """报告配置"""
    title: str = "数据分析报告"
    subtitle: str = ""
    author: str = "数据分析Agent"
    date: str = ""
    logo_path: Optional[str] = None
    template_path: Optional[str] = None
    output_format: str = "word"  # word 或 ppt
    
    def __post_init__(self):
        if not self.date:
            self.date = datetime.now().strftime("%Y年%m月%d日")


class ReportBuilder(ABC):
    """报告生成器基类"""
    
    def __init__(self, config: ReportConfig):
        self.config = config
    
    @abstractmethod
    def build(self, analysis_result: AnalysisResult, output_path: str) -> str:
        """生成报告"""
        pass
    
    @abstractmethod
    def add_title_page(self):
        """添加标题页"""
        pass
    
    @abstractmethod
    def add_summary_section(self, summary: str):
        """添加摘要章节"""
        pass
    
    @abstractmethod
    def add_findings_section(self, findings: List[str]):
        """添加关键发现章节"""
        pass
    
    @abstractmethod
    def add_table(self, table_data: Dict[str, Any]):
        """添加表格"""
        pass
    
    @abstractmethod
    def add_recommendations_section(self, recommendations: List[str]):
        """添加建议章节"""
        pass


class WordReportBuilder(ReportBuilder):
    """Word报告生成器"""
    
    def __init__(self, config: ReportConfig):
        super().__init__(config)
        self.doc = Document()
        self._setup_styles()
    
    def _setup_styles(self):
        """设置文档样式"""
        # 设置默认字体
        style = self.doc.styles['Normal']
        font = style.font
        font.name = '微软雅黑'
        font.size = Pt(11)
        
        # 标题样式
        for i in range(1, 4):
            style_name = f'Heading {i}'
            if style_name in self.doc.styles:
                heading_style = self.doc.styles[style_name]
                heading_style.font.name = '微软雅黑'
                heading_style.font.bold = True
                if i == 1:
                    heading_style.font.size = Pt(18)
                    heading_style.font.color.rgb = RGBColor(0, 51, 102)
                elif i == 2:
                    heading_style.font.size = Pt(14)
                    heading_style.font.color.rgb = RGBColor(0, 76, 153)
                else:
                    heading_style.font.size = Pt(12)
    
    def build(self, analysis_result: AnalysisResult, output_path: str) -> str:
        """生成Word报告"""
        # 添加标题页
        self.add_title_page()
        
        # 添加目录占位
        self.doc.add_paragraph()
        self.doc.add_paragraph("目录", style='Heading 1')
        self.doc.add_paragraph("（生成后请更新目录）")
        self.doc.add_page_break()
        
        # 添加摘要
        self.add_summary_section(analysis_result.summary)
        
        # 新增：添加任务分析
        if analysis_result.task_analysis:
            self.add_task_analysis_section(analysis_result.task_analysis)
        
        # 新增：添加问题总结
        if analysis_result.problem_summary:
            self.add_problem_summary_section(analysis_result.problem_summary)
        
        # 新增：添加经营目标
        if analysis_result.business_goals:
            self.add_business_goals_section(analysis_result.business_goals)
        
        # 新增：添加改善方法
        if analysis_result.improvement_methods:
            self.add_improvement_methods_section(analysis_result.improvement_methods)
        
        # 添加关键发现
        self.add_findings_section(analysis_result.key_findings)
        
        # 添加关键指标
        if analysis_result.metrics:
            self.add_metrics_section(analysis_result.metrics)
        
        # 添加数据表格
        for table_data in analysis_result.tables[:3]:  # 最多3个表格
            self.add_table(table_data)
        
        # 添加排名分析
        if analysis_result.rankings:
            self.add_rankings_section(analysis_result.rankings)
        
        # 添加趋势分析
        if analysis_result.trends:
            self.add_trends_section(analysis_result.trends)
        
        # 添加建议
        self.add_recommendations_section(analysis_result.recommendations)
        
        # 保存文档
        self.doc.save(output_path)
        return output_path
    
    def add_title_page(self):
        """添加标题页"""
        # 添加空行使标题居中
        for _ in range(5):
            self.doc.add_paragraph()
        
        # 主标题
        title = self.doc.add_paragraph()
        title_run = title.add_run(self.config.title)
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0, 51, 102)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 副标题
        if self.config.subtitle:
            subtitle = self.doc.add_paragraph()
            subtitle_run = subtitle.add_run(self.config.subtitle)
            subtitle_run.font.size = Pt(16)
            subtitle_run.font.color.rgb = RGBColor(102, 102, 102)
            subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 空行
        for _ in range(10):
            self.doc.add_paragraph()
        
        # 作者和日期
        info = self.doc.add_paragraph()
        info_run = info.add_run(f"生成者: {self.config.author}\n日期: {self.config.date}")
        info_run.font.size = Pt(12)
        info.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        self.doc.add_page_break()
    
    def add_summary_section(self, summary: str):
        """添加摘要章节"""
        self.doc.add_heading('一、报告摘要', level=1)
        
        para = self.doc.add_paragraph()
        para.add_run(summary)
        para.paragraph_format.first_line_indent = Cm(0.75)
        para.paragraph_format.line_spacing = 1.5
        
        self.doc.add_paragraph()
    
    def add_findings_section(self, findings: List[str]):
        """添加关键发现章节"""
        self.doc.add_heading('二、关键发现', level=1)
        
        for i, finding in enumerate(findings, 1):
            para = self.doc.add_paragraph()
            para.add_run(f"{i}. ").bold = True
            para.add_run(finding)
            para.paragraph_format.line_spacing = 1.5
        
        self.doc.add_paragraph()
    
    def add_metrics_section(self, metrics: Dict[str, Any]):
        """添加关键指标章节"""
        self.doc.add_heading('三、关键指标', level=1)
        
        # 创建指标表格
        table = self.doc.add_table(rows=1, cols=2)
        table.style = 'Table Grid'
        
        # 表头
        header_cells = table.rows[0].cells
        header_cells[0].text = '指标名称'
        header_cells[1].text = '指标值'
        
        for cell in header_cells:
            cell.paragraphs[0].runs[0].font.bold = True
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 数据行
        for name, value in metrics.items():
            row_cells = table.add_row().cells
            row_cells[0].text = str(name)
            row_cells[1].text = str(value)
        
        self.doc.add_paragraph()
    
    def add_table(self, table_data: Dict[str, Any]):
        """添加数据表格"""
        title = table_data.get('title', '数据表格')
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not headers or not rows:
            return
        
        # 限制列数和行数
        max_cols = min(len(headers), 8)
        max_rows = min(len(rows), 15)
        
        self.doc.add_heading(title, level=2)
        
        table = self.doc.add_table(rows=1, cols=max_cols)
        table.style = 'Table Grid'
        
        # 表头
        header_cells = table.rows[0].cells
        for i, header in enumerate(headers[:max_cols]):
            header_cells[i].text = str(header)[:20]  # 限制表头长度
            header_cells[i].paragraphs[0].runs[0].font.bold = True
            header_cells[i].paragraphs[0].runs[0].font.size = Pt(9)
        
        # 数据行
        for row in rows[:max_rows]:
            row_cells = table.add_row().cells
            for i, value in enumerate(row[:max_cols]):
                row_cells[i].text = str(value)[:30] if value else ""
                row_cells[i].paragraphs[0].runs[0].font.size = Pt(9)
        
        self.doc.add_paragraph()
    
    def add_rankings_section(self, rankings: List[Dict[str, Any]]):
        """添加排名分析章节"""
        self.doc.add_heading('四、排名分析', level=1)
        
        for ranking in rankings[:2]:  # 最多2个排名
            sheet = ranking.get('sheet', '')
            data = ranking.get('data', [])
            
            if data:
                self.doc.add_heading(f'{sheet} 排名', level=2)
                
                # 创建简单列表
                for i, item in enumerate(data[:10], 1):
                    para = self.doc.add_paragraph()
                    item_str = ", ".join(f"{k}: {v}" for k, v in item.items() if v)
                    para.add_run(f"{i}. {item_str}")
        
        self.doc.add_paragraph()
    
    def add_trends_section(self, trends: List[Dict[str, Any]]):
        """添加趋势分析章节"""
        self.doc.add_heading('五、趋势分析', level=1)
        
        for trend in trends[:3]:
            column = trend.get('column', '')
            values = trend.get('values', [])
            
            if values:
                para = self.doc.add_paragraph()
                para.add_run(f"• {column}: ").bold = True
                
                # 计算趋势
                if len(values) >= 2:
                    try:
                        numeric_values = [float(v) for v in values if v is not None]
                        if numeric_values:
                            avg = sum(numeric_values) / len(numeric_values)
                            trend_direction = "上升" if numeric_values[-1] > numeric_values[0] else "下降"
                            para.add_run(f"平均值 {avg:.2f}，整体呈{trend_direction}趋势")
                    except:
                        para.add_run(f"数据点数: {len(values)}")
        
        self.doc.add_paragraph()
    
    def add_recommendations_section(self, recommendations):
        """添加建议章节 - 支持详细的结构化建议"""
        self.doc.add_heading('六、建议与行动计划', level=1)
        
        for i, rec in enumerate(recommendations, 1):
            # 支持新的结构化建议格式
            if isinstance(rec, dict):
                category = rec.get('category', '')
                action = rec.get('action', '')
                target = rec.get('target', '')
                priority = rec.get('priority', '')
                
                # 类别标题
                para = self.doc.add_paragraph()
                para.add_run(f"【{category}】").bold = True
                para.paragraph_format.space_before = Pt(12)
                
                # 执行策略
                if action:
                    action_para = self.doc.add_paragraph()
                    action_para.add_run("执行策略：").bold = True
                    action_para.add_run(action)
                    action_para.paragraph_format.left_indent = Cm(0.5)
                
                # 量化目标
                if target:
                    target_para = self.doc.add_paragraph()
                    target_para.add_run("量化目标：").bold = True
                    target_run = target_para.add_run(target)
                    target_run.font.color.rgb = RGBColor(0, 102, 51)  # 绿色突出
                    target_para.paragraph_format.left_indent = Cm(0.5)
                
                # 优先级
                if priority:
                    priority_para = self.doc.add_paragraph()
                    priority_para.add_run("优先级：").bold = True
                    priority_para.add_run(priority)
                    priority_para.paragraph_format.left_indent = Cm(0.5)
            else:
                # 兼容旧的字符串格式
                para = self.doc.add_paragraph()
                para.add_run(f"建议{i}: ").bold = True
                para.add_run(str(rec))
                para.paragraph_format.line_spacing = 1.5
        
        self.doc.add_paragraph()
    
    def add_task_analysis_section(self, task_analysis: Dict[str, Any]):
        """添加任务分析章节"""
        self.doc.add_heading('二、任务分析', level=1)
        
        task_name = task_analysis.get('task_name', '')
        current_status = task_analysis.get('current_status', '')
        target_gap = task_analysis.get('target_gap', '')
        completion_rate = task_analysis.get('completion_rate', '')
        key_blockers = task_analysis.get('key_blockers', [])
        
        # 任务名称
        if task_name:
            para = self.doc.add_paragraph()
            para.add_run("任务名称：").bold = True
            para.add_run(task_name)
        
        # 当前状态
        if current_status:
            para = self.doc.add_paragraph()
            para.add_run("当前状态：").bold = True
            para.add_run(current_status)
        
        # 完成率
        if completion_rate:
            para = self.doc.add_paragraph()
            para.add_run("完成率：").bold = True
            run = para.add_run(str(completion_rate))
            run.font.color.rgb = RGBColor(0, 102, 51)
            run.font.bold = True
        
        # 与目标差距
        if target_gap:
            para = self.doc.add_paragraph()
            para.add_run("与目标差距：").bold = True
            run = para.add_run(target_gap)
            run.font.color.rgb = RGBColor(204, 0, 0)
        
        # 关键阻碍因素
        if key_blockers:
            para = self.doc.add_paragraph()
            para.add_run("关键阻碍因素：").bold = True
            for blocker in key_blockers:
                blocker_para = self.doc.add_paragraph()
                blocker_para.add_run(f"  • {blocker}")
                blocker_para.paragraph_format.left_indent = Cm(0.5)
        
        self.doc.add_paragraph()
    
    def add_problem_summary_section(self, problem_summary: List[Dict[str, Any]]):
        """添加问题总结章节"""
        self.doc.add_heading('三、问题总结', level=1)
        
        for i, problem in enumerate(problem_summary, 1):
            category = problem.get('category', '')
            problem_desc = problem.get('problem', '')
            current_value = problem.get('current_value', '')
            target_value = problem.get('target_value', '')
            gap = problem.get('gap', '')
            impact = problem.get('impact', '')
            root_cause = problem.get('root_cause', '')
            
            # 问题标题
            self.doc.add_heading(f"问题{i}：【{category}】{problem_desc}", level=2)
            
            # 创建问题详情表格
            table = self.doc.add_table(rows=5, cols=2)
            table.style = 'Table Grid'
            
            rows_data = [
                ("当前值", current_value),
                ("目标值", target_value),
                ("差距", gap),
                ("影响程度", impact),
                ("根因分析", root_cause)
            ]
            
            for row_idx, (label, value) in enumerate(rows_data):
                table.rows[row_idx].cells[0].text = label
                table.rows[row_idx].cells[0].paragraphs[0].runs[0].font.bold = True
                table.rows[row_idx].cells[1].text = str(value)
            
            self.doc.add_paragraph()
        
        self.doc.add_paragraph()
    
    def add_business_goals_section(self, business_goals: List[Dict[str, Any]]):
        """添加经营目标章节"""
        self.doc.add_heading('四、经营目标', level=1)
        
        # 创建目标汇总表格
        table = self.doc.add_table(rows=1, cols=5)
        table.style = 'Table Grid'
        
        # 表头
        headers = ["目标名称", "目标值", "当前值", "达成时间", "优先级"]
        for i, header in enumerate(headers):
            table.rows[0].cells[i].text = header
            table.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True
        
        # 数据行
        for goal in business_goals:
            row = table.add_row()
            row.cells[0].text = goal.get('goal_name', '')
            row.cells[1].text = str(goal.get('target_value', ''))
            row.cells[2].text = str(goal.get('current_value', ''))
            row.cells[3].text = goal.get('timeline', '')
            row.cells[4].text = goal.get('priority', '')
        
        self.doc.add_paragraph()
        
        # 详细说明
        for goal in business_goals:
            goal_name = goal.get('goal_name', '')
            rationale = goal.get('rationale', '')
            if rationale:
                para = self.doc.add_paragraph()
                para.add_run(f"• {goal_name}：").bold = True
                para.add_run(rationale)
        
        self.doc.add_paragraph()
    
    def add_improvement_methods_section(self, improvement_methods: List[Dict[str, Any]]):
        """添加改善方法章节"""
        self.doc.add_heading('五、改善方法', level=1)
        
        for i, method in enumerate(improvement_methods, 1):
            goal_ref = method.get('goal_ref', '')
            category = method.get('category', '')
            method_desc = method.get('method', '')
            action_steps = method.get('action_steps', [])
            responsible = method.get('responsible', '')
            resources_needed = method.get('resources_needed', '')
            expected_result = method.get('expected_result', '')
            timeline = method.get('timeline', '')
            
            # 改善方法标题
            self.doc.add_heading(f"方法{i}：【{category}】{goal_ref}", level=2)
            
            # 具体方法
            para = self.doc.add_paragraph()
            para.add_run("改善方法：").bold = True
            method_run = para.add_run(method_desc)
            method_run.font.color.rgb = RGBColor(0, 51, 102)
            
            # 执行步骤
            if action_steps:
                para = self.doc.add_paragraph()
                para.add_run("执行步骤：").bold = True
                for step_idx, step in enumerate(action_steps, 1):
                    step_para = self.doc.add_paragraph()
                    step_para.add_run(f"  {step_idx}. {step}")
                    step_para.paragraph_format.left_indent = Cm(0.5)
            
            # 责任方
            if responsible:
                para = self.doc.add_paragraph()
                para.add_run("责任方：").bold = True
                para.add_run(responsible)
            
            # 所需资源
            if resources_needed:
                para = self.doc.add_paragraph()
                para.add_run("所需资源：").bold = True
                para.add_run(resources_needed)
            
            # 预期效果
            if expected_result:
                para = self.doc.add_paragraph()
                para.add_run("预期效果：").bold = True
                result_run = para.add_run(expected_result)
                result_run.font.color.rgb = RGBColor(0, 102, 51)
            
            # 执行时间
            if timeline:
                para = self.doc.add_paragraph()
                para.add_run("执行时间：").bold = True
                para.add_run(timeline)
            
            self.doc.add_paragraph()
        
        self.doc.add_paragraph()


class PPTReportBuilder(ReportBuilder):
    """PPT报告生成器"""
    
    def __init__(self, config: ReportConfig):
        super().__init__(config)
        self.prs = Presentation()
        # 设置幻灯片大小为16:9
        self.prs.slide_width = PptInches(13.333)
        self.prs.slide_height = PptInches(7.5)
    
    def build(self, analysis_result: AnalysisResult, output_path: str) -> str:
        """生成PPT报告"""
        # 添加标题页
        self.add_title_page()
        
        # 添加摘要页
        self.add_summary_section(analysis_result.summary)
        
        # 新增：添加任务分析页
        if analysis_result.task_analysis:
            self.add_task_analysis_slide(analysis_result.task_analysis)
        
        # 新增：添加问题总结页
        if analysis_result.problem_summary:
            self.add_problem_summary_slides(analysis_result.problem_summary)
        
        # 新增：添加经营目标页
        if analysis_result.business_goals:
            self.add_business_goals_slide(analysis_result.business_goals)
        
        # 新增：添加改善方法页
        if analysis_result.improvement_methods:
            self.add_improvement_methods_slides(analysis_result.improvement_methods)
        
        # 添加关键发现页
        self.add_findings_section(analysis_result.key_findings)
        
        # 添加关键指标页
        if analysis_result.metrics:
            self.add_metrics_slide(analysis_result.metrics)
        
        # 添加数据表格页
        for table_data in analysis_result.tables[:2]:  # 最多2个表格
            self.add_table(table_data)
        
        # 添加建议页
        self.add_recommendations_section(analysis_result.recommendations)
        
        # 添加结束页
        self.add_end_slide()
        
        # 保存
        self.prs.save(output_path)
        return output_path
    
    def _add_slide(self, layout_index: int = 6):
        """添加幻灯片"""
        # layout_index 6 通常是空白布局
        try:
            layout = self.prs.slide_layouts[layout_index]
        except:
            layout = self.prs.slide_layouts[0]
        return self.prs.slides.add_slide(layout)
    
    def add_title_page(self):
        """添加标题页"""
        slide = self._add_slide(0)  # 标题布局
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(2.5), 
            PptInches(12.333), PptInches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = self.config.title
        title_para.font.size = PptPt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PptRGBColor(0, 51, 102)
        title_para.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        if self.config.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                PptInches(0.5), PptInches(4), 
                PptInches(12.333), PptInches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = self.config.subtitle
            subtitle_para.font.size = PptPt(24)
            subtitle_para.font.color.rgb = PptRGBColor(102, 102, 102)
            subtitle_para.alignment = PP_ALIGN.CENTER
        
        # 添加日期
        date_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(6), 
            PptInches(12.333), PptInches(0.5)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = self.config.date
        date_para.font.size = PptPt(16)
        date_para.alignment = PP_ALIGN.CENTER
    
    def add_summary_section(self, summary: str):
        """添加摘要页"""
        slide = self._add_slide()
        
        # 标题
        self._add_slide_title(slide, "报告摘要")
        
        # 内容
        content_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(1.5), 
            PptInches(12.333), PptInches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        para = content_frame.paragraphs[0]
        para.text = summary
        para.font.size = PptPt(18)
        para.line_spacing = 1.5
    
    def add_findings_section(self, findings: List[str]):
        """添加关键发现页"""
        slide = self._add_slide()
        
        # 标题
        self._add_slide_title(slide, "关键发现")
        
        # 内容
        content_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(1.5), 
            PptInches(12.333), PptInches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, finding in enumerate(findings):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            para.text = f"• {finding}"
            para.font.size = PptPt(16)
            para.space_after = PptPt(12)
    
    def add_metrics_slide(self, metrics: Dict[str, Any]):
        """添加关键指标页"""
        slide = self._add_slide()
        
        # 标题
        self._add_slide_title(slide, "关键指标")
        
        # 创建指标卡片布局
        items = list(metrics.items())
        cols = min(3, len(items))
        rows = (len(items) + cols - 1) // cols
        
        card_width = 3.5
        card_height = 1.5
        start_x = (13.333 - cols * card_width - (cols - 1) * 0.3) / 2
        start_y = 2
        
        for i, (name, value) in enumerate(items[:6]):  # 最多6个指标
            col = i % cols
            row = i // cols
            
            x = start_x + col * (card_width + 0.3)
            y = start_y + row * (card_height + 0.3)
            
            # 添加卡片背景
            shape = slide.shapes.add_shape(
                1,  # 矩形
                PptInches(x), PptInches(y),
                PptInches(card_width), PptInches(card_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGBColor(240, 240, 240)
            shape.line.color.rgb = PptRGBColor(200, 200, 200)
            
            # 添加指标名称
            name_box = slide.shapes.add_textbox(
                PptInches(x + 0.1), PptInches(y + 0.2),
                PptInches(card_width - 0.2), PptInches(0.5)
            )
            name_para = name_box.text_frame.paragraphs[0]
            name_para.text = str(name)
            name_para.font.size = PptPt(12)
            name_para.font.color.rgb = PptRGBColor(102, 102, 102)
            
            # 添加指标值
            value_box = slide.shapes.add_textbox(
                PptInches(x + 0.1), PptInches(y + 0.7),
                PptInches(card_width - 0.2), PptInches(0.6)
            )
            value_para = value_box.text_frame.paragraphs[0]
            value_para.text = str(value)
            value_para.font.size = PptPt(24)
            value_para.font.bold = True
            value_para.font.color.rgb = PptRGBColor(0, 51, 102)
    
    def add_table(self, table_data: Dict[str, Any]):
        """添加表格页"""
        title = table_data.get('title', '数据表格')
        headers = table_data.get('headers', [])
        rows = table_data.get('rows', [])
        
        if not headers or not rows:
            return
        
        slide = self._add_slide()
        
        # 标题
        self._add_slide_title(slide, title)
        
        # 限制表格大小
        max_cols = min(len(headers), 6)
        max_rows = min(len(rows), 8)
        
        # 计算表格位置和大小
        table_width = PptInches(12)
        table_height = PptInches(4.5)
        left = PptInches(0.667)
        top = PptInches(1.8)
        
        # 创建表格
        table = slide.shapes.add_table(
            max_rows + 1, max_cols,
            left, top, table_width, table_height
        ).table
        
        # 设置表头
        for i, header in enumerate(headers[:max_cols]):
            cell = table.cell(0, i)
            cell.text = str(header)[:15]
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGBColor(0, 51, 102)
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = PptRGBColor(255, 255, 255)
            para.font.size = PptPt(10)
            para.font.bold = True
        
        # 填充数据
        for row_idx, row in enumerate(rows[:max_rows]):
            for col_idx, value in enumerate(row[:max_cols]):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)[:20] if value else ""
                para = cell.text_frame.paragraphs[0]
                para.font.size = PptPt(9)
    
    def add_recommendations_section(self, recommendations):
        """添加建议页 - 支持详细的结构化建议，每个建议一页"""
        # 检查是否是结构化建议
        has_structured = any(isinstance(rec, dict) for rec in recommendations)
        
        if has_structured:
            # 结构化建议：每个建议一页
            for rec in recommendations:
                if isinstance(rec, dict):
                    self._add_structured_recommendation_slide(rec)
                else:
                    self._add_simple_recommendation_slide([rec])
        else:
            # 简单建议：一页展示所有
            self._add_simple_recommendation_slide(recommendations)
    
    def _add_structured_recommendation_slide(self, rec: dict):
        """添加结构化建议页"""
        slide = self._add_slide()
        
        category = rec.get('category', '建议')
        action = rec.get('action', '')
        target = rec.get('target', '')
        priority = rec.get('priority', '')
        
        # 标题
        self._add_slide_title(slide, f"【{category}】")
        
        # 执行策略
        y_pos = 1.8
        if action:
            action_box = slide.shapes.add_textbox(
                PptInches(0.5), PptInches(y_pos),
                PptInches(12.333), PptInches(1.2)
            )
            action_frame = action_box.text_frame
            action_frame.word_wrap = True
            
            label_para = action_frame.paragraphs[0]
            label_para.text = "执行策略"
            label_para.font.size = PptPt(14)
            label_para.font.bold = True
            label_para.font.color.rgb = PptRGBColor(102, 102, 102)
            
            content_para = action_frame.add_paragraph()
            content_para.text = action
            content_para.font.size = PptPt(20)
            content_para.font.bold = True
            y_pos += 1.5
        
        # 量化目标 - 突出显示
        if target:
            target_shape = slide.shapes.add_shape(
                1, PptInches(0.5), PptInches(y_pos),
                PptInches(12.333), PptInches(1.2)
            )
            target_shape.fill.solid()
            target_shape.fill.fore_color.rgb = PptRGBColor(0, 102, 51)
            target_shape.line.fill.background()
            
            target_box = slide.shapes.add_textbox(
                PptInches(0.7), PptInches(y_pos + 0.15),
                PptInches(12), PptInches(1)
            )
            target_frame = target_box.text_frame
            target_frame.word_wrap = True
            
            label_para = target_frame.paragraphs[0]
            label_para.text = "量化目标"
            label_para.font.size = PptPt(12)
            label_para.font.color.rgb = PptRGBColor(255, 255, 255)
            
            content_para = target_frame.add_paragraph()
            content_para.text = target
            content_para.font.size = PptPt(24)
            content_para.font.bold = True
            content_para.font.color.rgb = PptRGBColor(255, 255, 255)
            y_pos += 1.6
        
        # 优先级
        if priority:
            priority_box = slide.shapes.add_textbox(
                PptInches(0.5), PptInches(y_pos),
                PptInches(12.333), PptInches(1)
            )
            priority_frame = priority_box.text_frame
            priority_frame.word_wrap = True
            
            label_para = priority_frame.paragraphs[0]
            label_para.text = "优先级"
            label_para.font.size = PptPt(14)
            label_para.font.bold = True
            label_para.font.color.rgb = PptRGBColor(102, 102, 102)
            
            content_para = priority_frame.add_paragraph()
            content_para.text = priority
            content_para.font.size = PptPt(18)
            content_para.font.color.rgb = PptRGBColor(204, 102, 0)
    
    def _add_simple_recommendation_slide(self, recommendations):
        """添加简单建议页"""
        slide = self._add_slide()
        
        # 标题
        self._add_slide_title(slide, "建议与行动计划")
        
        # 内容
        content_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(1.5), 
            PptInches(12.333), PptInches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, rec in enumerate(recommendations):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            para.text = f"{i+1}. {rec}"
            para.font.size = PptPt(16)
            para.space_after = PptPt(12)
    
    def add_task_analysis_slide(self, task_analysis: Dict[str, Any]):
        """添加任务分析页"""
        slide = self._add_slide()
        self._add_slide_title(slide, "任务分析")
        
        task_name = task_analysis.get('task_name', '')
        current_status = task_analysis.get('current_status', '')
        target_gap = task_analysis.get('target_gap', '')
        completion_rate = task_analysis.get('completion_rate', '')
        key_blockers = task_analysis.get('key_blockers', [])
        
        y_pos = 1.5
        
        # 任务名称
        if task_name:
            box = slide.shapes.add_textbox(PptInches(0.5), PptInches(y_pos), PptInches(12.333), PptInches(0.6))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = f"任务：{task_name}"
            para.font.size = PptPt(20)
            para.font.bold = True
            y_pos += 0.7
        
        # 完成率 - 大字突出
        if completion_rate:
            shape = slide.shapes.add_shape(1, PptInches(0.5), PptInches(y_pos), PptInches(4), PptInches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGBColor(0, 102, 51)
            shape.line.fill.background()
            
            box = slide.shapes.add_textbox(PptInches(0.7), PptInches(y_pos + 0.1), PptInches(3.6), PptInches(1))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = "完成率"
            para.font.size = PptPt(12)
            para.font.color.rgb = PptRGBColor(255, 255, 255)
            para2 = frame.add_paragraph()
            para2.text = str(completion_rate)
            para2.font.size = PptPt(36)
            para2.font.bold = True
            para2.font.color.rgb = PptRGBColor(255, 255, 255)
        
        # 差距 - 红色突出
        if target_gap:
            shape = slide.shapes.add_shape(1, PptInches(5), PptInches(y_pos), PptInches(7.833), PptInches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGBColor(204, 0, 0)
            shape.line.fill.background()
            
            box = slide.shapes.add_textbox(PptInches(5.2), PptInches(y_pos + 0.1), PptInches(7.4), PptInches(1))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = "与目标差距"
            para.font.size = PptPt(12)
            para.font.color.rgb = PptRGBColor(255, 255, 255)
            para2 = frame.add_paragraph()
            para2.text = target_gap
            para2.font.size = PptPt(20)
            para2.font.bold = True
            para2.font.color.rgb = PptRGBColor(255, 255, 255)
        
        y_pos += 1.5
        
        # 当前状态
        if current_status:
            box = slide.shapes.add_textbox(PptInches(0.5), PptInches(y_pos), PptInches(12.333), PptInches(0.6))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = f"当前状态：{current_status}"
            para.font.size = PptPt(16)
            y_pos += 0.8
        
        # 关键阻碍因素
        if key_blockers:
            box = slide.shapes.add_textbox(PptInches(0.5), PptInches(y_pos), PptInches(12.333), PptInches(3))
            frame = box.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.text = "关键阻碍因素："
            para.font.size = PptPt(14)
            para.font.bold = True
            for blocker in key_blockers[:4]:
                p = frame.add_paragraph()
                p.text = f"• {blocker}"
                p.font.size = PptPt(14)
    
    def add_problem_summary_slides(self, problem_summary: List[Dict[str, Any]]):
        """添加问题总结页"""
        slide = self._add_slide()
        self._add_slide_title(slide, "问题总结")
        
        # 创建问题汇总表格
        problems = problem_summary[:5]  # 最多5个问题
        if not problems:
            return
        
        table = slide.shapes.add_table(len(problems) + 1, 4, PptInches(0.5), PptInches(1.5), PptInches(12.333), PptInches(5)).table
        
        # 表头
        headers = ["问题类别", "问题描述", "差距", "影响"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = PptRGBColor(0, 51, 102)
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = PptRGBColor(255, 255, 255)
            para.font.size = PptPt(12)
            para.font.bold = True
        
        # 数据行
        for row_idx, problem in enumerate(problems):
            table.cell(row_idx + 1, 0).text = problem.get('category', '')
            table.cell(row_idx + 1, 1).text = problem.get('problem', '')[:30]
            table.cell(row_idx + 1, 2).text = problem.get('gap', '')
            table.cell(row_idx + 1, 3).text = problem.get('impact', '')
            for col in range(4):
                table.cell(row_idx + 1, col).text_frame.paragraphs[0].font.size = PptPt(10)
    
    def add_business_goals_slide(self, business_goals: List[Dict[str, Any]]):
        """添加经营目标页"""
        slide = self._add_slide()
        self._add_slide_title(slide, "经营目标")
        
        goals = business_goals[:5]
        if not goals:
            return
        
        # 创建目标卡片
        card_width = 2.4
        card_height = 2.5
        start_x = 0.5
        start_y = 1.5
        
        for i, goal in enumerate(goals):
            col = i % 5
            x = start_x + col * (card_width + 0.2)
            
            # 卡片背景
            shape = slide.shapes.add_shape(1, PptInches(x), PptInches(start_y), PptInches(card_width), PptInches(card_height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PptRGBColor(240, 240, 240)
            shape.line.color.rgb = PptRGBColor(200, 200, 200)
            
            # 目标名称
            box = slide.shapes.add_textbox(PptInches(x + 0.1), PptInches(start_y + 0.1), PptInches(card_width - 0.2), PptInches(0.5))
            para = box.text_frame.paragraphs[0]
            para.text = goal.get('goal_name', '')[:12]
            para.font.size = PptPt(11)
            para.font.bold = True
            para.font.color.rgb = PptRGBColor(0, 51, 102)
            
            # 目标值
            box = slide.shapes.add_textbox(PptInches(x + 0.1), PptInches(start_y + 0.6), PptInches(card_width - 0.2), PptInches(0.8))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = "目标"
            para.font.size = PptPt(9)
            para.font.color.rgb = PptRGBColor(102, 102, 102)
            para2 = frame.add_paragraph()
            para2.text = str(goal.get('target_value', ''))
            para2.font.size = PptPt(18)
            para2.font.bold = True
            para2.font.color.rgb = PptRGBColor(0, 102, 51)
            
            # 当前值
            box = slide.shapes.add_textbox(PptInches(x + 0.1), PptInches(start_y + 1.5), PptInches(card_width - 0.2), PptInches(0.6))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = f"当前：{goal.get('current_value', '')}"
            para.font.size = PptPt(10)
            
            # 优先级
            priority = goal.get('priority', '')
            box = slide.shapes.add_textbox(PptInches(x + 0.1), PptInches(start_y + 2.1), PptInches(card_width - 0.2), PptInches(0.3))
            para = box.text_frame.paragraphs[0]
            para.text = f"{priority} | {goal.get('timeline', '')}"
            para.font.size = PptPt(9)
            para.font.color.rgb = PptRGBColor(204, 102, 0) if priority == 'P0' else PptRGBColor(102, 102, 102)
    
    def add_improvement_methods_slides(self, improvement_methods: List[Dict[str, Any]]):
        """添加改善方法页 - 每个方法一页"""
        for method in improvement_methods[:5]:  # 最多5个方法
            slide = self._add_slide()
            
            goal_ref = method.get('goal_ref', '')
            category = method.get('category', '')
            method_desc = method.get('method', '')
            action_steps = method.get('action_steps', [])
            expected_result = method.get('expected_result', '')
            timeline = method.get('timeline', '')
            
            self._add_slide_title(slide, f"【{category}】{goal_ref}")
            
            y_pos = 1.5
            
            # 改善方法
            box = slide.shapes.add_textbox(PptInches(0.5), PptInches(y_pos), PptInches(12.333), PptInches(1))
            frame = box.text_frame
            frame.word_wrap = True
            para = frame.paragraphs[0]
            para.text = "改善方法"
            para.font.size = PptPt(12)
            para.font.color.rgb = PptRGBColor(102, 102, 102)
            para2 = frame.add_paragraph()
            para2.text = method_desc
            para2.font.size = PptPt(18)
            para2.font.bold = True
            y_pos += 1.2
            
            # 执行步骤
            if action_steps:
                box = slide.shapes.add_textbox(PptInches(0.5), PptInches(y_pos), PptInches(12.333), PptInches(2))
                frame = box.text_frame
                frame.word_wrap = True
                para = frame.paragraphs[0]
                para.text = "执行步骤："
                para.font.size = PptPt(12)
                para.font.bold = True
                for step_idx, step in enumerate(action_steps[:3], 1):
                    p = frame.add_paragraph()
                    p.text = f"{step_idx}. {step}"
                    p.font.size = PptPt(14)
                y_pos += 1.8
            
            # 预期效果 - 绿色背景
            if expected_result:
                shape = slide.shapes.add_shape(1, PptInches(0.5), PptInches(y_pos), PptInches(8), PptInches(0.8))
                shape.fill.solid()
                shape.fill.fore_color.rgb = PptRGBColor(0, 102, 51)
                shape.line.fill.background()
                
                box = slide.shapes.add_textbox(PptInches(0.7), PptInches(y_pos + 0.1), PptInches(7.6), PptInches(0.6))
                frame = box.text_frame
                para = frame.paragraphs[0]
                para.text = f"预期效果：{expected_result}"
                para.font.size = PptPt(14)
                para.font.bold = True
                para.font.color.rgb = PptRGBColor(255, 255, 255)
            
            # 执行时间
            if timeline:
                box = slide.shapes.add_textbox(PptInches(9), PptInches(y_pos), PptInches(4), PptInches(0.8))
                para = box.text_frame.paragraphs[0]
                para.text = f"时间：{timeline}"
                para.font.size = PptPt(12)
                para.font.color.rgb = PptRGBColor(102, 102, 102)
    
    def add_end_slide(self):
        """添加结束页"""
        slide = self._add_slide()
        
        # 感谢文字
        thanks_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(3), 
            PptInches(12.333), PptInches(1.5)
        )
        thanks_frame = thanks_box.text_frame
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.text = "谢谢！"
        thanks_para.font.size = PptPt(44)
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = PptRGBColor(0, 51, 102)
        thanks_para.alignment = PP_ALIGN.CENTER
        
        # 生成信息
        info_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(5), 
            PptInches(12.333), PptInches(0.5)
        )
        info_frame = info_box.text_frame
        info_para = info_frame.paragraphs[0]
        info_para.text = f"由数据分析Agent自动生成 | {self.config.date}"
        info_para.font.size = PptPt(12)
        info_para.font.color.rgb = PptRGBColor(150, 150, 150)
        info_para.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide, title: str):
        """添加幻灯片标题"""
        title_box = slide.shapes.add_textbox(
            PptInches(0.5), PptInches(0.3), 
            PptInches(12.333), PptInches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = PptPt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = PptRGBColor(0, 51, 102)


def create_report(
    analysis_result: AnalysisResult,
    output_path: str,
    config: ReportConfig = None,
    format: str = "word"
) -> str:
    """
    创建报告
    
    Args:
        analysis_result: 分析结果
        output_path: 输出路径
        config: 报告配置
        format: 输出格式 (word/ppt)
        
    Returns:
        str: 生成的报告文件路径
    """
    if config is None:
        config = ReportConfig()
    
    config.output_format = format
    
    if format == "ppt":
        builder = PPTReportBuilder(config)
        if not output_path.endswith('.pptx'):
            output_path = output_path.rsplit('.', 1)[0] + '.pptx'
    else:
        builder = WordReportBuilder(config)
        if not output_path.endswith('.docx'):
            output_path = output_path.rsplit('.', 1)[0] + '.docx'
    
    return builder.build(analysis_result, output_path)
